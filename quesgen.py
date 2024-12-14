import random
import spacy
import pickle
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
from pdfplumber import open as pdf_open
from docx import Document
from pptx import Presentation

# Load spaCy NER model
nlp = spacy.load("en_core_web_sm")

# Function to load the saved model and tokenizer from pickle files in the pklmodels directory
def load_model_from_pkl(model_name):
    model_path = f"pklmodels/{model_name}_model.pkl"
    tokenizer_path = f"pklmodels/{model_name}_tokenizer.pkl"
    
    with open(model_path, "rb") as model_file:
        model = pickle.load(model_file)
    with open(tokenizer_path, "rb") as tokenizer_file:
        tokenizer = pickle.load(tokenizer_file)
    
    print(f"{model_name} and tokenizer loaded successfully.")
    return model, tokenizer

# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    with pdf_open(pdf_path) as pdf:
        text = """ """
        for page in pdf.pages:
            text += page.extract_text()
    return text

# Function to extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = """ """
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = "" ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to split text into chunks for processing
def split_text_into_chunks(text, chunk_size=1024):
    chunks = []
    text_split = text.split("\n")
    chunk = ""
    for line in text_split:
        if len(chunk + line) > chunk_size:
            chunks.append(chunk)
            chunk = line
        else:
            chunk += line + "\n"
    if chunk:
        chunks.append(chunk)  # Add last chunk if any
    return chunks

# Named Entity Recognition function to extract entities
def extract_entities(text):
    doc = nlp(text)
    entities = {ent.label_: ent.text for ent in doc.ents}
    return entities

# Function to generate a question from the context using FLAN-T5 or your saved model
def generate_question_from_context(model, tokenizer, context):
    prompt = f"Context: {context}\nGenerate a relevant question:"
    inputs = tokenizer(prompt, return_tensors="pt", padding=True, truncation=True)
    outputs = model.generate(inputs["input_ids"], max_length=50, num_beams=5)
    generated_question = tokenizer.decode(outputs[0], skip_special_tokens=True)
    return generated_question

# Function to extract the correct answer using QA model
def extract_correct_answer(context, question):
    # Use the question-answering pipeline
    nlp_qa = pipeline("question-answering", model="distilbert-base-uncased-distilled-squad", tokenizer="distilbert-base-uncased-distilled-squad")
    result = nlp_qa(question=question, context=context)
    return result['answer']

# Function to generate multiple choice questions based on context and entities
def generate_mcq_questions_from_context(model, tokenizer, text):
    chunks = split_text_into_chunks(text)
    score = 0  # To track the score
    question_answers = []  # To track the user's answers
    
    # Loop through chunks and generate questions for each
    for chunk in chunks:
        # Generate a question from the current chunk
        question = generate_question_from_context(model, tokenizer, chunk)
        
        # Extract the correct answer using the QA model
        correct_answer = extract_correct_answer(chunk, question)
        
        # Generate distractors (incorrect answers) using context-based approach
        options = [correct_answer]

        # Extract random non-entity words from the chunk (ensuring correct answer is included)
        non_entity_words = [word for word in chunk.split() if word not in correct_answer]
        random.shuffle(non_entity_words)
        
        # Add 3 random non-entity distractors to the options
        distractors = non_entity_words[:3]
        options.extend(distractors)

        # Shuffle the options to mix correct answer with distractors
        random.shuffle(options)
        
        print(f"\nQuestion: {question}")
        print("Options: ")
        for idx, option in enumerate(options, 1):
            print(f"{idx}. {option}")
        
        # Ask the user to choose an answer
        while True:
            try:
                user_answer_idx = int(input("\nEnter the number corresponding to your answer (1-4), or 0 to skip the question: "))
                if user_answer_idx not in [0, 1, 2, 3, 4]:
                    print("Please choose a valid option between 1 and 4, or 0 to skip.")
                    continue
                user_answer = options[user_answer_idx - 1] if user_answer_idx != 0 else None
                break
            except ValueError:
                print("Please enter a valid number.")

        # Store the answer (None if skipped)
        question_answers.append({
            'question': question,
            'user_answer': user_answer,
            'correct_answer': correct_answer
        })

        if user_answer == correct_answer:
            print("Correct!")
            score += 1  # Increment score for correct answer
        elif user_answer is not None:
            print(f"Incorrect. The correct answer is: {correct_answer}")
    
    return score, question_answers

# Load the saved models and tokenizer
model, tokenizer = load_model_from_pkl("flan_t5")  # Update with your model's name

# Initialize score
score = 0

# Example usage:
# Load a PDF, DOCX, or PPTX file and process it
text = extract_text_from_pdf("example.pdf")  # For PDF
# text = extract_text_from_docx("example.docx")  # For DOCX
# text = extract_text_from_pptx("example.pptx")  # For PPTX

# Ask the user if they want to attempt the quiz
attempt_quiz = input("Do you want to attempt the quiz? (y/n): ").strip().lower()

if attempt_quiz == 'y':
    while True:
        score, question_answers = generate_mcq_questions_from_context(model, tokenizer, text)
        print(f"\nYour current score is: {score}")
        next_action = input("\nDo you want to continue to the next question? (y/n): ").lower()
        if next_action != 'y':
            print(f"Final score: {score}. Exiting quiz. Goodbye!")
            break
    
    # Display the final answers after exiting the quiz
    print("\nQuiz Results:")
    for answer in question_answers:
        print(f"Question: {answer['question']}")
        print(f"Your Answer: {answer['user_answer']}")
        print(f"Correct Answer: {answer['correct_answer']}")
        print()
else:
    print("Quiz skipped. Goodbye!")
